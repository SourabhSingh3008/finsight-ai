from docx import Document
from docx.shared import Pt, Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
import os

def create_docx():
    doc = Document()
    doc.add_heading('FinSight AI - Project Summary', 0)

    doc.add_heading('1. Executive Summary', level=1)
    doc.add_paragraph(
        "FinSight AI is an intelligent Expense & Subscription Analyzer web application designed to help users take control "
        "of their personal finances. Built as a pure client-side application using HTML, CSS (glassmorphism UI), and Vanilla JS, "
        "it utilizes the Gemini 1.5 Flash API to process unstructured bank transaction text. The application categorizes expenses, "
        "detects active and hidden subscriptions, and acts as a personalized financial advisor by generating actionable insights."
    )

    doc.add_heading('2. Problem Addressed', level=1)
    doc.add_paragraph(
        "Messy and unstructured bank statements make it difficult for individuals to understand their spending habits or "
        "keep track of recurring subscriptions. Hidden costs add up quickly, and traditional budgeting apps require tedious manual categorization or "
        "unsafe bank integrations."
    )

    doc.add_heading('3. The Solution', level=1)
    doc.add_paragraph(
        "FinSight AI solves this by employing generative AI to parse raw text (directly copied from banking portals or PDFs) "
        "into structured financial data. It automatically identifies categories and subscriptions and calculates total spend, all while "
        "ensuring data is processed securely through standard API calls without storing data on a backend server."
    )

    doc.add_heading('4. Technical Stack', level=1)
    doc.add_paragraph(
        "- Frontend Core: HTML5, CSS3, Vanilla JavaScript\n"
        "- Aesthetics: Dark mode, Glassmorphism gradients, Micro-animations\n"
        "- AI Integration: Directly uses Google's Gemini 1.5 Flash API via REST fetch requests\n"
        "- Architecture: Client-side only; users supply their own Gemini API key (stored securely in localstorage), eliminating backend infrastructure needs."
    )

    doc.add_heading('5. Key Features', level=1)
    doc.add_paragraph(
        "- Unstructured Data Import: No CSV formatting required; users just paste raw text.\n"
        "- Automated Categorization: The AI infers and groups spending dynamically.\n"
        "- Subscription Detection: Identifies recurring payments (like Netflix or Uber One) and sums the monthly flow.\n"
        "- Financial Advisor: Generates three custom, actionable insights based heavily on the user's specific expenditure patterns."
    )

    doc.add_page_break()
    doc.save('FinSight_AI_Summary.docx')
    print("DOCX created successfully.")

def create_pptx():
    prs = Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "FinSight AI"
    subtitle.text = "Intelligent Expense & Subscription Analyzer\nMBA Startup Sprint Project"

    # Slide 2: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "The Problem"
    tf = content.text_frame
    tf.text = "1. Unstructured Banking Data: Bank statements are messy and hard to read."
    p = tf.add_paragraph()
    p.text = "2. Trial Subscriptions: Users lose money to forgotten recurring subscriptions."
    p = tf.add_paragraph()
    p.text = "3. Generic Advice: Existing apps fail to provide personalized, actionable financial insights."

    # Slide 3: The Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Our Solution: FinSight AI"
    tf = content.text_frame
    tf.text = "A sleek, web-based tool powered by Gemini AI."
    p = tf.add_paragraph()
    p.text = "Process: User pastes raw transaction text -> AI extracts Data -> Interactive Dashboard presents insights."
    p = tf.add_paragraph()
    p.text = "Features:\n- Auto-categorization of expenses.\n- Subscription tracking and alerts.\n- AI-generated personalized saving tips."

    # Slide 4: AI Usage
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "AI Integration & Stack"
    tf = content.text_frame
    tf.text = "We use GenAI as a core feature, evaluating context from raw financial strings."
    p = tf.add_paragraph()
    p.text = "- Model: Gemini 1.5 Flash for high-speed, deterministic JSON extraction."
    p = tf.add_paragraph()
    p.text = "- Prompt Engineering: Instructed mapping to JSON schema (total_spend, categories, insights)."
    p = tf.add_paragraph()
    p.text = "- Tech Stack: Pure Client-Side architecture (HTML/CSS/Vanilla JS) enabling secure API key handling locally."

    # Slide 5: Live Demo & Impact
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Live Demo & Future Scope"
    tf = content.text_frame
    tf.text = "It just works! Simply paste transactions and receive an instant financial health check."
    p = tf.add_paragraph()
    p.text = "Future Scope:"
    p = tf.add_paragraph()
    p.text = "- Support for uploading bank PDFs directly."
    p = tf.add_paragraph()
    p.text = "- Integration with standard financial APIs (e.g. Plaid)."

    prs.save('FinSight_AI_Pitch_Deck.pptx')
    print("PPTX created successfully.")

if __name__ == '__main__':
    create_docx()
    create_pptx()
